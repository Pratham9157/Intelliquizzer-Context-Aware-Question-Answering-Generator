"""
MODULE 1 — Text Extractor
==========================
Handles extraction of raw text from:
  • Plain text  (pasted / .txt)
  • PDF         (pdfplumber — layout-aware, column-safe)
  • PowerPoint  (python-pptx — slide-by-slide, preserves hierarchy)

Returns the raw combined text string; cleaning is done in Module 2.

Design notes:
  - PDF extraction uses word-level sorting to handle multi-column layouts
  - PPTX extraction tries to preserve logical reading order per slide
  - All extractors strip binary noise but keep paragraph structure (double newlines)
"""

from __future__ import annotations

import io
import re
import logging
from pathlib import Path
from typing import Tuple

logger = logging.getLogger(__name__)

# ── optional dependencies ─────────────────────────────────────────────────────
try:
    import pdfplumber
    _PDF_OK = True
except ImportError:
    _PDF_OK = False
    logger.warning("pdfplumber not installed — PDF extraction disabled.")

try:
    from pptx import Presentation
    from pptx.util import Pt
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False
    logger.warning("python-pptx not installed — PPTX extraction disabled.")


# ─────────────────────────────────────────────────────────────────────────────
# Internal helpers
# ─────────────────────────────────────────────────────────────────────────────

def _remove_markdown_links(text: str) -> str:
    """Strip [[n](url)] citation fragments that appear in web-scraped text."""
    text = re.sub(r'\[\[?\d+\]?\([^)]*\)\]?', '', text)
    text = re.sub(r'\[https?://[^\]]+\]', '', text)
    text = re.sub(r'https?://\S+', '', text)
    return text


def _normalise_whitespace(text: str) -> str:
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


# ─────────────────────────────────────────────────────────────────────────────
# Extractors
# ─────────────────────────────────────────────────────────────────────────────

def from_text(raw: str) -> str:
    """Process pasted / plain text."""
    raw = _remove_markdown_links(raw)
    return _normalise_whitespace(raw)


def from_pdf(data: bytes) -> str:
    """
    Extract text from PDF bytes using pdfplumber.

    Strategy:
      1. Use extract_words() with sorted positions so multi-column PDFs
         are read in natural left-right, top-bottom order.
      2. Fall back to extract_text() per page if word extraction is empty.
      3. Pages separated by double newlines for paragraph detection downstream.
    """
    if not _PDF_OK:
        raise ImportError("pdfplumber is required for PDF extraction. "
                          "pip install pdfplumber")
    pages: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            # Attempt word-level extraction (handles multi-column)
            words = page.extract_words(
                x_tolerance=3, y_tolerance=3,
                keep_blank_chars=False, use_text_flow=True
            )
            if words:
                # Sort top-to-bottom, left-to-right
                words.sort(key=lambda w: (round(w['top'] / 10), w['x0']))
                page_text = ' '.join(w['text'] for w in words)
            else:
                page_text = page.extract_text() or ''
            if page_text.strip():
                pages.append(page_text.strip())

    return _normalise_whitespace('\n\n'.join(pages))


def from_pptx(data: bytes) -> str:
    """
    Extract text from PPTX bytes using python-pptx.

    Per slide:
      • Title shape extracted first (if present)
      • Body shapes follow in top-to-bottom order
      • Bullet hierarchy preserved via indentation → kept as plain text
    Slides separated by double newlines.
    """
    if not _PPTX_OK:
        raise ImportError("python-pptx is required for PPTX extraction. "
                          "pip install python-pptx")
    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        title_text = ''

        # Sort shapes top-to-bottom
        shapes = sorted(slide.shapes, key=lambda s: (s.top or 0))

        for shape in shapes:
            if not shape.has_text_frame:
                continue
            # Detect title placeholder
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title_text = shape.text_frame.text.strip()
            else:
                lines = []
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        lines.append(line)
                if lines:
                    parts.append('\n'.join(lines))

        slide_parts: list[str] = []
        if title_text:
            slide_parts.append(f'[Slide {slide_idx}: {title_text}]')
        slide_parts.extend(parts)

        combined = '\n'.join(slide_parts).strip()
        if combined:
            slides.append(combined)

    return _normalise_whitespace('\n\n'.join(slides))


def extract(source) -> str:
    """
    Auto-detect source type and extract text.

    Args:
        source: One of:
          - str                     → treated as raw text
          - bytes with magic bytes  → auto-detected PDF/PPTX
          - Streamlit UploadedFile  → detected via .name attribute

    Returns:
        Extracted raw text string.
    """
    # Streamlit UploadedFile
    if hasattr(source, 'name') and hasattr(source, 'read'):
        name = source.name.lower()
        data = source.read()
        if name.endswith('.pdf'):
            return from_pdf(data)
        elif name.endswith(('.pptx', '.ppt')):
            return from_pptx(data)
        else:
            return from_text(data.decode('utf-8', errors='ignore'))

    # Plain string
    if isinstance(source, str):
        return from_text(source)

    # Raw bytes — sniff magic bytes
    if isinstance(source, bytes):
        if source[:4] == b'%PDF':
            return from_pdf(source)
        if source[:2] == b'PK':          # ZIP-based (DOCX/PPTX)
            return from_pptx(source)
        return from_text(source.decode('utf-8', errors='ignore'))

    raise TypeError(f"Unsupported source type: {type(source)}")
