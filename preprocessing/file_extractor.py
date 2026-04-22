"""
MODULE 1: Text Extraction Layer
================================
Handles extraction of text from multiple file formats:
- PDF files (PyMuPDF/fitz)
- PowerPoint presentations (python-pptx)
- Plain text files

Purpose:
    Convert any input document format to clean, plain text string
    suitable for downstream NLP processing.

Key Features:
    - Robust error handling
    - Encoding detection for text files
    - Metadata extraction and logging
    - File type validation
    - Clean text normalization
"""

import os
import re
from pathlib import Path
from typing import Optional, Dict, Tuple
import logging

# Auto-setup dependencies on first import
try:
    from setup import ensure_dependencies
    ensure_dependencies()
except ImportError:
    pass  # setup.py not available, assume manual setup

# ============================================================================
# CONFIGURATION CONSTANTS
# ============================================================================

# Text separators
PAGE_SEPARATOR = "\n\n--- PAGE BREAK ---\n\n"
SLIDE_SEPARATOR = "\n\n--- SLIDE BREAK ---\n\n"

# Logging configuration
DEFAULT_LOG_LEVEL = logging.INFO

# Supported file formats
SUPPORTED_FORMATS = {'.pdf', '.pptx', '.ppt', '.txt'}

# Text encoding fallback order
ENCODING_FALLBACK = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-1']

# Configure logging
logger = logging.getLogger(__name__)


class TextExtractor:
    """
    Main class for extracting text from various document formats.
    
    Supported formats:
        - PDF: .pdf (requires PyMuPDF)
        - PowerPoint: .pptx, .ppt (requires python-pptx)
        - Text: .txt
    
    Example:
        >>> extractor = TextExtractor()
        >>> text = extractor.extract("path/to/file.pdf")
        >>> print(len(text), "characters extracted")
        
        >>> # With debug logging
        >>> import logging
        >>> extractor = TextExtractor(log_level=logging.DEBUG)
    """
    
    def __init__(self, log_level: int = DEFAULT_LOG_LEVEL):
        """
        Initialize the TextExtractor and check dependencies.
        
        Args:
            log_level: Logging level (logging.DEBUG, logging.INFO, etc.)
                      Default: logging.INFO
        """
        # Configure logging for this module
        if not logger.handlers:
            handler = logging.StreamHandler()
            formatter = logging.Formatter(
                '%(asctime)s - %(name)s - %(levelname)s - %(message)s'
            )
            handler.setFormatter(formatter)
            logger.addHandler(handler)
        
        logger.setLevel(log_level)
        
        self.supported_formats = SUPPORTED_FORMATS
        self._check_dependencies()
    
    def _check_dependencies(self) -> None:
        """Check if required libraries are installed."""
        self.pdf_available = self._check_pdf_lib()
        self.pptx_available = self._check_pptx_lib()
        
        if not self.pdf_available:
            logger.warning("PyMuPDF not installed. PDF extraction unavailable. "
                          "Install with: pip install PyMuPDF")
        
        if not self.pptx_available:
            logger.warning("python-pptx not installed. PowerPoint extraction unavailable. "
                          "Install with: pip install python-pptx")
    
    @staticmethod
    def _check_pdf_lib() -> bool:
        """Check if PyMuPDF is available."""
        try:
            import fitz  # noqa
            return True
        except ImportError:
            return False
    
    @staticmethod
    def _check_pptx_lib() -> bool:
        """Check if python-pptx is available."""
        try:
            from pptx import Presentation  # noqa
            return True
        except ImportError:
            return False
    
    def extract(self, file_path: str) -> str:
        """
        Main entry point: extract text from any supported file format.
        
        Args:
            file_path: Path to the input file
            
        Returns:
            Extracted and cleaned text string
            
        Raises:
            FileNotFoundError: If file does not exist
            ValueError: If file format is not supported or dependencies missing
            Exception: For PDF/PPTX parsing errors
        """
        file_path = Path(file_path)
        
        # Validate file existence
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Get file extension
        file_ext = file_path.suffix.lower()
        
        # Validate format support
        if file_ext not in self.supported_formats:
            raise ValueError(
                f"Unsupported file format: {file_ext}. "
                f"Supported: {self.supported_formats}"
            )
        
        # Route to appropriate extractor
        logger.info(f"Extracting text from {file_path.name} ({file_ext})")
        
        try:
            if file_ext == '.pdf':
                text = self._extract_pdf(str(file_path))
            elif file_ext in {'.pptx', '.ppt'}:
                text = self._extract_pptx(str(file_path))
            elif file_ext == '.txt':
                text = self._extract_txt(str(file_path))
            else:
                raise ValueError(f"Unexpected file extension: {file_ext}")
            
            # Clean and normalize
            text = self._normalize_text(text)
            
            logger.info(f"Successfully extracted {len(text)} characters")
            return text
            
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise
    
    def _extract_pdf(self, file_path: str) -> str:
        """
        Extract text from PDF file using PyMuPDF.
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text from all pages
            
        Raises:
            ImportError: If PyMuPDF not installed
            Exception: If PDF parsing fails
        """
        if not self.pdf_available:
            raise ImportError(
                "PyMuPDF not installed. Install with: pip install PyMuPDF"
            )
        
        import fitz  # Lazy import to avoid hard dependency
        
        text_parts = []
        
        try:
            # Open PDF document
            doc = fitz.open(file_path)
            logger.info(f"PDF has {doc.page_count} pages")
            
            # Extract text from each page
            for page_num in range(doc.page_count):
                try:
                    page = doc[page_num]
                    page_text = page.get_text()
                    text_parts.append(page_text)
                    logger.debug(f"Extracted page {page_num + 1}/{doc.page_count}")
                except Exception as e:
                    logger.warning(f"Failed to extract page {page_num + 1}: {str(e)}")
                    continue
            
            doc.close()
            
        except fitz.FileError as e:
            raise Exception(f"Failed to open PDF file: {str(e)}")
        except Exception as e:
            raise Exception(f"PDF extraction error: {str(e)}")
        
        # Join all pages with page separator
        return PAGE_SEPARATOR.join(text_parts)
    
    def _extract_pptx(self, file_path: str) -> str:
        """
        Extract text from PowerPoint presentation.
        
        Extracts text from:
        - Slide titles
        - Slide text boxes
        - Speaker notes
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Extracted text from all slides
            
        Raises:
            ImportError: If python-pptx not installed
            Exception: If PPTX parsing fails
        """
        if not self.pptx_available:
            raise ImportError(
                "python-pptx not installed. Install with: pip install python-pptx"
            )
        
        from pptx import Presentation  # Lazy import
        
        text_parts = []
        
        try:
            # Load presentation
            prs = Presentation(file_path)
            logger.info(f"Presentation has {len(prs.slides)} slides")
            
            # Extract from each slide
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extract from shapes (text boxes, titles, etc.)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                # Extract speaker notes if available
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_text.append(f"[Speaker Notes: {notes_text}]")
                
                if slide_text:
                    text_parts.append("\n".join(slide_text))
                    logger.debug(f"Extracted slide {slide_num + 1}/{len(prs.slides)}")
            
        except Exception as e:
            raise Exception(f"PPTX extraction error: {str(e)}")
        
        # Join all slides with slide separator
        return SLIDE_SEPARATOR.join(text_parts)
    
    def _extract_txt(self, file_path: str) -> str:
        """
        Extract text from plain text file.
        
        Handles various encodings with fallback strategy:
        1. Try UTF-8 (most common)
        2. Try UTF-8 with error handling
        3. Try common encodings (latin-1, cp1252)
        
        Args:
            file_path: Path to text file
            
        Returns:
            Text content of file
            
        Raises:
            Exception: If unable to decode file
        """
        for encoding in ENCODING_FALLBACK:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                logger.info(f"Successfully decoded with {encoding}")
                return text
            except (UnicodeDecodeError, LookupError) as e:
                logger.debug(f"Failed with {encoding}: {str(e)}")
                continue
        
        # If all encodings fail, read with error replacement
        logger.warning("All standard encodings failed, reading with error replacement")
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                text = f.read()
            logger.warning("File read with replacement characters for undecodable bytes")
            return text
        except Exception as e:
            raise Exception(f"Unable to read text file: {str(e)}")
    
    @staticmethod
    def _normalize_text(text: str) -> str:
        """
        Normalize extracted text for downstream processing.
        
        Operations:
        - Remove multiple consecutive whitespace
        - Normalize line breaks
        - Remove control characters (except newlines)
        - Strip leading/trailing whitespace
        
        Args:
            text: Raw extracted text
            
        Returns:
            Normalized text
        """
        # Remove control characters except tab, newline, carriage return
        text = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', text)
        
        # Normalize line breaks (multiple -> single)
        text = re.sub(r'\r\n', '\n', text)  # Windows line breaks
        text = re.sub(r'\r', '\n', text)    # Old Mac line breaks
        text = re.sub(r'\n\n+', '\n\n', text)  # Multiple newlines -> double newline
        
        # Normalize spaces (multiple -> single, but preserve double newlines)
        lines = text.split('\n\n')
        normalized_lines = [
            re.sub(r'[ \t]+', ' ', line).strip()
            for line in lines
        ]
        text = '\n\n'.join(normalized_lines)
        
        # Final strip
        text = text.strip()
        
        return text
    
    def extract_with_metadata(self, file_path: str) -> Dict:
        """
        Extract text and return metadata about the extraction.
        
        Useful for logging and debugging.
        
        Returns:
            Dictionary containing:
            - 'text': Extracted text
            - 'file_name': Original file name
            - 'file_size_kb': File size in KB
            - 'file_format': File extension
            - 'char_count': Number of characters
            - 'line_count': Number of lines
        """
        file_path = Path(file_path)
        
        text = self.extract(str(file_path))
        
        return {
            'text': text,
            'file_name': file_path.name,
            'file_size_kb': round(file_path.stat().st_size / 1024, 2),
            'file_format': file_path.suffix,
            'char_count': len(text),
            'line_count': text.count('\n') + 1
        }


# ============================================================================
# TEST FUNCTION
# ============================================================================

def test_file_extractor():
    """
    Test the TextExtractor with sample files.
    
    This function demonstrates:
    - Creating a TextExtractor instance
    - Extracting from different file types
    - Handling errors gracefully
    - Using metadata extraction
    """
    print("\n" + "="*70)
    print("TEXT EXTRACTION MODULE - TEST SUITE")
    print("="*70 + "\n")
    
    extractor = TextExtractor()
    
    # Test 1: Create sample text file
    print("[TEST 1] Creating and extracting sample TXT file...")
    sample_txt_path = "sample_test.txt"
    sample_txt_content = """
    Artificial Intelligence and Machine Learning
    
    Artificial Intelligence (AI) is the simulation of human intelligence
    by computer systems. These systems can learn from experience and improve
    their performance over time without being explicitly programmed.
    
    Machine Learning is a subset of AI that focuses on teaching computers
    to learn from data and make predictions without explicit programming.
    """
    
    with open(sample_txt_path, 'w', encoding='utf-8') as f:
        f.write(sample_txt_content)
    
    try:
        text = extractor.extract(sample_txt_path)
        print(f"✓ Successfully extracted {len(text)} characters from TXT file")
        print(f"  First 100 chars: {text[:100]}...\n")
    except Exception as e:
        print(f"✗ Error: {e}\n")
    finally:
        os.remove(sample_txt_path)
    
    # Test 2: File not found error handling
    print("[TEST 2] Testing file not found error handling...")
    try:
        extractor.extract("nonexistent_file.pdf")
        print("✗ Should have raised FileNotFoundError\n")
    except FileNotFoundError as e:
        print(f"✓ Correctly raised FileNotFoundError: {e}\n")
    
    # Test 3: Unsupported format error handling
    print("[TEST 3] Testing unsupported file format error handling...")
    
    # Create sample unsupported file
    sample_docx_path = "sample_test.docx"
    with open(sample_docx_path, 'w') as f:
        f.write("dummy")
    
    try:
        extractor.extract(sample_docx_path)
        print("✗ Should have raised ValueError\n")
    except ValueError as e:
        print(f"✓ Correctly raised ValueError: {e}\n")
    finally:
        os.remove(sample_docx_path)
    
    # Test 4: Metadata extraction
    print("[TEST 4] Testing metadata extraction...")
    sample_txt_path = "sample_metadata_test.txt"
    sample_content = "Test content for metadata extraction.\nLine 2.\nLine 3."
    
    with open(sample_txt_path, 'w', encoding='utf-8') as f:
        f.write(sample_content)
    
    try:
        metadata = extractor.extract_with_metadata(sample_txt_path)
        print(f"✓ Metadata extraction successful:")
        print(f"  - File name: {metadata['file_name']}")
        print(f"  - File size: {metadata['file_size_kb']} KB")
        print(f"  - Characters: {metadata['char_count']}")
        print(f"  - Lines: {metadata['line_count']}\n")
    except Exception as e:
        print(f"✗ Error: {e}\n")
    finally:
        os.remove(sample_txt_path)
    
    # Test 5: Text normalization
    print("[TEST 5] Testing text normalization...")
    sample_txt_path = "sample_normalization_test.txt"
    messy_content = "Text   with\n\nmultiple    spaces\r\nand weird\n\n\nline breaks"
    
    with open(sample_txt_path, 'w', encoding='utf-8') as f:
        f.write(messy_content)
    
    try:
        normalized_text = extractor.extract(sample_txt_path)
        print(f"✓ Text normalization successful:")
        print(f"  Original length: {len(messy_content)}")
        print(f"  Normalized length: {len(normalized_text)}")
        print(f"  Normalized text:\n    {repr(normalized_text)}\n")
    except Exception as e:
        print(f"✗ Error: {e}\n")
    finally:
        os.remove(sample_txt_path)
    
    # Test 6: Encoding detection (UTF-8 with BOM)
    print("[TEST 6] Testing encoding detection (UTF-8 with BOM)...")
    sample_txt_path = "sample_encoding_test.txt"
    
    with open(sample_txt_path, 'w', encoding='utf-8-sig') as f:
        f.write("Text with UTF-8 BOM encoding")
    
    try:
        text = extractor.extract(sample_txt_path)
        print(f"✓ Successfully handled UTF-8 BOM encoding")
        print(f"  Extracted text: {text}\n")
    except Exception as e:
        print(f"✗ Error: {e}\n")
    finally:
        os.remove(sample_txt_path)
    
    print("="*70)
    print("TEST SUITE COMPLETED")
    print("="*70 + "\n")


if __name__ == "__main__":
    test_file_extractor()
